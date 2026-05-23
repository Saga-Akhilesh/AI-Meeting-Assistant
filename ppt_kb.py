import os
import json
import re
from pathlib import Path
from pptx import Presentation

def ensure_dir(p):
    Path(p).mkdir(parents=True, exist_ok=True)

def clean_text(s: str) -> str:
    if not s:
        return ""
    s = re.sub(r"\s+", " ", s)
    return s.strip()

def export_slides_to_images(pptx_path: str, out_dir: str):
    """
    Uses PowerPoint COM automation to export slides as PNG.
    Requires: MS PowerPoint + pywin32
    """
    import win32com.client

    ensure_dir(out_dir)
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = 1

    pres = ppt.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
    pres.Export(os.path.abspath(out_dir), "PNG")  # Slide1.PNG, Slide2.PNG...
    pres.Close()
    ppt.Quit()

def extract_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            t = clean_text(shape.text)
            if t:
                parts.append(t)
    return "\n".join(parts)

def extract_speaker_notes(prs: Presentation, slide_index: int) -> str:
    try:
        notes_slide = prs.slides[slide_index].notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            return clean_text(notes_slide.notes_text_frame.text)
    except Exception:
        pass
    return ""

def build_slide_context(slide_no: int, ppt_text: str, ocr_text: str, notes: str) -> str:
    return f"""
SLIDE {slide_no}

[PPT TEXT]
{ppt_text}

[OCR TEXT]
{ocr_text}

[SPEAKER NOTES]
{notes}
""".strip()

def build_kb(pptx_path: str, slides_png_dir: str, kb_out_path: str, ocr_func):
    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        ppt_text = extract_slide_text(slide)
        notes = extract_speaker_notes(prs, i - 1)

        img_path = os.path.join(slides_png_dir, f"Slide{i}.PNG")

        ocr_text = ""
        if os.path.exists(img_path):
            ocr_text = ocr_func(img_path)

        context = build_slide_context(i, ppt_text, ocr_text, notes)

        slides.append({
            "slide_no": i,
            "ppt_text": ppt_text,
            "notes": notes,
            "ocr_text": ocr_text,
            "image_path": img_path,
            "context": context
        })

    ensure_dir(os.path.dirname(kb_out_path))
    with open(kb_out_path, "w", encoding="utf-8") as f:
        json.dump(slides, f, indent=2, ensure_ascii=False)

    return slides
